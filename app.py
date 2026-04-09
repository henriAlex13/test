import streamlit as st
import pandas as pd
import sqlite3
import hashlib
import plotly.express as px
import plotly.graph_objects as go
from datetime import datetime
import io
import os

st.set_page_config(
    page_title="Dashboard KPI – SocGen CI",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

DB_PATH = "kpi_socgen.db"


# ── Utilitaires ────────────────────────────────────────────────────────────────
def _hash(p: str) -> str:
    return hashlib.sha256(p.encode()).hexdigest()


# ── Utilisateurs ───────────────────────────────────────────────────────────────
USERS = {
    "admin":          {"password": _hash("Admin@2026"),       "role": "admin", "nom_ref": None,             "display": "Administrateur"},
    "dg":             {"password": _hash("DG@2026"),          "role": "dg",    "nom_ref": None,             "display": "Directeur Général"},
    "jacques":        {"password": _hash("Jacques@2026"),     "role": "user",  "nom_ref": "jacques",        "display": "Jacques Akponi"},
    "yssouf":         {"password": _hash("Yssouf@2026"),      "role": "user",  "nom_ref": "yssouf",         "display": "Yssouf Soumahoro"},
    "hafsatou":       {"password": _hash("Hafsatou@2026"),    "role": "user",  "nom_ref": "Hafsatou",       "display": "Hafsatou Thiam"},
    "aymard":         {"password": _hash("Aymard@2026"),      "role": "user",  "nom_ref": "aymard",         "display": "Aymard Konan"},
    "jean-joseph":    {"password": _hash("JeanJoseph@2026"),  "role": "user",  "nom_ref": "Jean-Joseph",    "display": "Jean-Joseph Kouassi"},
    "esther":         {"password": _hash("Esther@2026"),      "role": "user",  "nom_ref": "esther",         "display": "Esther Toure"},
    "nicanor":        {"password": _hash("Nicanor@2026"),     "role": "user",  "nom_ref": "Nicanor",        "display": "Nicanor Andju"},
    "isabelle":       {"password": _hash("Isabelle@2026"),    "role": "user",  "nom_ref": "Isabelle",       "display": "Isabelle Dirabou"},
    "serge-francois": {"password": _hash("Serge@2026"),       "role": "user",  "nom_ref": "Serge-francois", "display": "Serge-François Koffi"},
}

ROLE_LABELS = {"admin": "Admin", "dg": "DG", "user": "Responsable"}


# ── Base de données ────────────────────────────────────────────────────────────
def get_conn():
    return sqlite3.connect(DB_PATH, check_same_thread=False)


def init_db():
    conn = get_conn()
    c = conn.cursor()
    # Table KPI référentiel (dynamique)
    c.execute("""
        CREATE TABLE IF NOT EXISTS kpi_ref (
            id          INTEGER PRIMARY KEY AUTOINCREMENT,
            categorie   TEXT NOT NULL,
            kpi         TEXT NOT NULL,
            periodicite TEXT NOT NULL DEFAULT 'Mensuelle',
            email       TEXT,
            nom         TEXT,
            actif       INTEGER DEFAULT 1
        )
    """)
    # Table saisies
    c.execute("""
        CREATE TABLE IF NOT EXISTS saisies (
            id          INTEGER PRIMARY KEY AUTOINCREMENT,
            login       TEXT NOT NULL,
            nom_resp    TEXT NOT NULL,
            categorie   TEXT NOT NULL,
            kpi         TEXT NOT NULL,
            periodicite TEXT,
            periode     TEXT NOT NULL,
            valeur      REAL NOT NULL,
            unite       TEXT,
            commentaire TEXT,
            date_saisie TEXT NOT NULL
        )
    """)
    # Table logs
    c.execute("""
        CREATE TABLE IF NOT EXISTS logs_connexion (
            id        INTEGER PRIMARY KEY AUTOINCREMENT,
            login     TEXT,
            action    TEXT,
            timestamp TEXT
        )
    """)
    conn.commit()
    # Seed référentiel si vide
    if c.execute("SELECT COUNT(*) FROM kpi_ref").fetchone()[0] == 0:
        seed = [
            ("Activite Commerciale","Encours de depots retail","Mensuelle","jacques.akponi@socgen.com","jacques"),
            ("Activite Commerciale","Encours de depots corporate","Mensuelle","jacques.akponi@socgen.com","jacques"),
            ("Activite Commerciale","Taux de credit corporate","Mensuelle","yssouf.soumahoro@socgen.com","yssouf"),
            ("Activite Commerciale","Taux de credit Retail","Mensuelle","yssouf.soumahoro@socgen.com","yssouf"),
            ("Activite Commerciale","Taux de dossiers de credits echus Corporate","Mensuelle","Hafsatou.THIAM@socgen.com","Hafsatou"),
            ("Activite Commerciale","Taux de dossiers de credits echus Retail","Mensuelle","aymard.konan@socgen.com","aymard"),
            ("Activite Commerciale","Nombre de forcage Corporate","Mensuelle","Hafsatou.THIAM@socgen.com","Hafsatou"),
            ("Activite Commerciale","Nombre de Decaissement Non Conforme Corporate","Mensuelle","Hafsatou.THIAM@socgen.com","Hafsatou"),
            ("Performance Financiere","Commissions sur les produits Retail","Mensuelle","yssouf.soumahoro@socgen.com","yssouf"),
            ("Performance Financiere","Commissions sur les produits Corporate","Mensuelle","yssouf.soumahoro@socgen.com","yssouf"),
            ("Performance Financiere","Suspens Comptable","Trimestrielle","Jean-Joseph.KOUASSI@socgen.com","Jean-Joseph"),
            ("Qualite de portefeuille et risque","Nombre de PDI retail","Mensuelle","aymard.konan@socgen.com","aymard"),
            ("Qualite de portefeuille et risque","CNR global","Mensuelle","esther.toure@socgen.com","esther"),
            ("Qualite de portefeuille et risque","CNR retail","Mensuelle","esther.toure@socgen.com","esther"),
            ("Qualite de portefeuille et risque","CNR Corporate","Mensuelle","esther.toure@socgen.com","esther"),
            ("Ressources humaines","Taux de chaises vides par direction","Mensuelle","Nicanor.ANDJU@socgen.com","Nicanor"),
            ("Ressources humaines","Taux d'absenteisme par direction","Mensuelle","Nicanor.ANDJU@socgen.com","Nicanor"),
            ("Ressources humaines","Taux de turn over par direction","Mensuelle","Nicanor.ANDJU@socgen.com","Nicanor"),
            ("Gouvernance et Conformite","Nombre d'incidents financiers eleves et tres eleves","Mensuelle","Isabelle.Dirabou@socgen.com","Isabelle"),
            ("Gouvernance et Conformite","Nombre d'incidents operationnels","Mensuelle","Isabelle.Dirabou@socgen.com","Isabelle"),
            ("Gouvernance et Conformite","Montant des incidents Operationnels","Mensuelle","Isabelle.Dirabou@socgen.com","Isabelle"),
            ("Digitalisation","Disponibilite des GAB (Gobla et Top 5)","Mensuelle","Serge-francois.KOFFI@socgen.com","Serge-francois"),
        ]
        c.executemany(
            "INSERT INTO kpi_ref (categorie,kpi,periodicite,email,nom) VALUES (?,?,?,?,?)", seed
        )
        conn.commit()
    conn.close()


def load_ref() -> pd.DataFrame:
    conn = get_conn()
    df = pd.read_sql("SELECT * FROM kpi_ref WHERE actif=1 ORDER BY categorie,kpi", conn)
    conn.close()
    return df


def add_kpi_ref(categorie, kpi, periodicite, email, nom):
    conn = get_conn()
    conn.execute(
        "INSERT INTO kpi_ref (categorie,kpi,periodicite,email,nom) VALUES (?,?,?,?,?)",
        (categorie, kpi, periodicite, email, nom)
    )
    conn.commit()
    conn.close()


def delete_kpi_ref(kpi_id: int):
    conn = get_conn()
    conn.execute("UPDATE kpi_ref SET actif=0 WHERE id=?", (kpi_id,))
    conn.commit()
    conn.close()


def log_connexion(login, action):
    conn = get_conn()
    conn.execute(
        "INSERT INTO logs_connexion (login,action,timestamp) VALUES (?,?,?)",
        (login, action, datetime.now().strftime("%Y-%m-%d %H:%M:%S"))
    )
    conn.commit()
    conn.close()


def insert_saisie(data: dict):
    conn = get_conn()
    conn.execute("""
        INSERT INTO saisies
          (login,nom_resp,categorie,kpi,periodicite,periode,valeur,unite,commentaire,date_saisie)
        VALUES (?,?,?,?,?,?,?,?,?,?)
    """, (
        data["login"], data["nom_resp"], data["categorie"], data["kpi"],
        data["periodicite"], data["periode"], data["valeur"],
        data["unite"], data["commentaire"],
        datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    ))
    conn.commit()
    conn.close()


def load_saisies(login=None) -> pd.DataFrame:
    conn = get_conn()
    if login and login not in ("admin", "dg"):
        df = pd.read_sql(
            "SELECT * FROM saisies WHERE login=? ORDER BY date_saisie DESC",
            conn, params=(login,)
        )
    else:
        df = pd.read_sql("SELECT * FROM saisies ORDER BY date_saisie DESC", conn)
    conn.close()
    return df


def delete_saisie(row_id, login):
    conn = get_conn()
    if login in ("admin", "dg"):
        conn.execute("DELETE FROM saisies WHERE id=?", (row_id,))
    else:
        conn.execute("DELETE FROM saisies WHERE id=? AND login=?", (row_id, login))
    conn.commit()
    conn.close()


def load_logs() -> pd.DataFrame:
    conn = get_conn()
    df = pd.read_sql("SELECT * FROM logs_connexion ORDER BY timestamp DESC LIMIT 300", conn)
    conn.close()
    return df


init_db()


# ── Export XLSX formaté ────────────────────────────────────────────────────────
def to_excel_bytes(df: pd.DataFrame) -> bytes:
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="xlsxwriter") as writer:
        df.to_excel(writer, index=False, sheet_name="Historique KPI")
        wb  = writer.book
        ws  = writer.sheets["Historique KPI"]
        # Formats
        hdr_fmt = wb.add_format({"bold": True, "bg_color": "#CC0000", "font_color": "white",
                                  "border": 1, "align": "center", "valign": "vcenter"})
        cell_fmt = wb.add_format({"border": 1, "valign": "vcenter"})
        num_fmt  = wb.add_format({"border": 1, "num_format": "#,##0.0000", "valign": "vcenter"})
        date_fmt = wb.add_format({"border": 1, "valign": "vcenter", "bg_color": "#FFF5F5"})
        # Headers
        for col_num, col_name in enumerate(df.columns):
            ws.write(0, col_num, col_name, hdr_fmt)
        # Rows
        for row_num, row in enumerate(df.itertuples(index=False), start=1):
            for col_num, val in enumerate(row):
                col_name = df.columns[col_num]
                if col_name in ("valeur", "Valeur"):
                    ws.write(row_num, col_num, val, num_fmt)
                elif col_name in ("date_saisie", "Date saisie"):
                    ws.write(row_num, col_num, str(val), date_fmt)
                else:
                    ws.write(row_num, col_num, str(val) if pd.notna(val) else "", cell_fmt)
        # Largeur colonnes
        col_widths = {"id":"ID","login":"login","nom_resp":"Responsable",
                      "categorie":"Catégorie","kpi":"KPI","periodicite":"Périodicité",
                      "periode":"Période","valeur":"Valeur","unite":"Unité",
                      "commentaire":"Commentaire","date_saisie":"Date saisie"}
        widths = [5,12,20,28,45,15,15,12,10,35,22]
        for i, w in enumerate(widths[:len(df.columns)]):
            ws.set_column(i, i, w)
        ws.freeze_panes(1, 0)
    return buf.getvalue()


# ── Export PPT Dashboard ───────────────────────────────────────────────────────
def build_pptx(df_all: pd.DataFrame, df_ref: pd.DataFrame) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    RED   = RGBColor(0xCC, 0x00, 0x00)
    DARK  = RGBColor(0x1A, 0x1A, 0x2E)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT = RGBColor(0xF5, 0xF6, 0xFA)
    GRAY  = RGBColor(0x88, 0x88, 0x88)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank

    def add_rect(slide, l, t, w, h, fill=None, line=None):
        shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shape.line.fill.background() if line is None else None
        if line is None:
            shape.line.fill.background()
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
        else:
            shape.fill.background()
        return shape

    def add_text(slide, text, l, t, w, h, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf  = txb.text_frame
        tf.word_wrap = wrap
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return txb

    # ── SLIDE 1 : Titre ──────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    add_rect(slide, 0, 0, 0.5, 7.5, fill=RED)
    add_text(slide, "DASHBOARD KPI", 1.0, 1.8, 11, 1.2, size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "SocGen Côte d'Ivoire", 1.0, 3.0, 8, 0.7, size=22, bold=False, color=RGBColor(0xCC,0xCC,0xCC))
    add_text(slide, f"Rapport généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')}", 1.0, 4.0, 8, 0.5, size=14, color=GRAY)
    add_rect(slide, 1.0, 5.2, 2.5, 0.06, fill=RED)
    add_text(slide, "Confidentiel – Usage interne", 1.0, 5.4, 8, 0.4, size=11, color=GRAY)

    # ── Stats globales ────────────────────────────────────────────────────────
    total_kpis   = len(df_ref)
    kpis_saisis  = df_all["kpi"].nunique() if not df_all.empty else 0
    tx_compl     = round(kpis_saisis / total_kpis * 100) if total_kpis else 0
    nb_saisies   = len(df_all)
    nb_resps     = df_all["nom_resp"].nunique() if not df_all.empty else 0

    # ── SLIDE 2 : Métriques clés ─────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 1.0, fill=RED)
    add_text(slide, "VUE D'ENSEMBLE", 0.3, 0.18, 12, 0.65, size=22, bold=True, color=WHITE)
    add_text(slide, f"Rapport du {datetime.now().strftime('%B %Y')}", 9.0, 0.22, 4, 0.5, size=12, color=WHITE, align=PP_ALIGN.RIGHT)

    metrics = [
        (str(total_kpis),  "KPIs au référentiel", 0.4),
        (str(kpis_saisis), "KPIs saisis",          3.6),
        (f"{tx_compl}%",   "Taux de complétion",   6.8),
        (str(nb_saisies),  "Total saisies",        10.0),
    ]
    for val, lbl, x in metrics:
        add_rect(slide, x, 1.3, 2.8, 2.5, fill=LIGHT)
        add_rect(slide, x, 1.3, 0.12, 2.5, fill=RED)
        add_text(slide, val, x+0.25, 1.55, 2.5, 1.1, size=42, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, x+0.25, 2.65, 2.5, 0.5, size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # Complétion par catégorie
    if not df_all.empty:
        cat_t  = df_ref.groupby("categorie")["kpi"].count().reset_index(name="total")
        cat_s  = df_all.groupby("categorie")["kpi"].nunique().reset_index(name="saisis")
        cat_m  = cat_t.merge(cat_s, on="categorie", how="left").fillna(0)
        cat_m["pct"] = (cat_m["saisis"] / cat_m["total"] * 100).round(0)

        add_text(slide, "COMPLÉTION PAR CATÉGORIE", 0.4, 4.1, 12, 0.4, size=11, bold=True, color=RED)
        y = 4.6
        bar_max_w = 10.0
        for _, row in cat_m.iterrows():
            pct = row["pct"]
            add_text(slide, row["categorie"][:38], 0.4, y, 5.5, 0.32, size=9, color=DARK)
            add_rect(slide, 5.8, y+0.04, bar_max_w * pct / 100, 0.22, fill=RED)
            add_rect(slide, 5.8, y+0.04, bar_max_w, 0.22)  # contour gris
            add_text(slide, f"{int(pct)}%", 5.8 + bar_max_w * pct / 100 + 0.1, y, 0.8, 0.3, size=9, bold=True, color=DARK)
            y += 0.36

    # ── SLIDE 3 : Saisies par responsable ───────────────────────────────────
    if not df_all.empty:
        slide = prs.slides.add_slide(blank)
        add_rect(slide, 0, 0, 13.33, 1.0, fill=DARK)
        add_text(slide, "SAISIES PAR RESPONSABLE", 0.3, 0.18, 12, 0.65, size=22, bold=True, color=WHITE)

        by_r = df_all.groupby("nom_resp").size().reset_index(name="nb").sort_values("nb", ascending=False)
        colors = [RED, RGBColor(0xE8,0x36,0x36), RGBColor(0xFF,0x80,0x80),
                  RGBColor(0xFF,0xB3,0xB3), RGBColor(0x1A,0x1A,0x2E),
                  RGBColor(0x44,0x44,0x6E), RGBColor(0x88,0x88,0xAA)]

        y = 1.3
        col = 0
        for i, (_, row) in enumerate(by_r.iterrows()):
            x = 0.4 + col * 6.4
            clr = colors[i % len(colors)]
            add_rect(slide, x, y, 5.9, 1.4, fill=LIGHT)
            add_rect(slide, x, y, 0.12, 1.4, fill=clr)
            add_text(slide, str(row["nb"]), x+0.3, y+0.15, 2, 0.8, size=36, bold=True, color=clr)
            add_text(slide, row["nom_resp"], x+0.3, y+0.9, 5.4, 0.4, size=11, color=DARK)
            col += 1
            if col == 2:
                col = 0
                y += 1.6

    # ── SLIDE 4 : Tableau historique (20 dernières lignes) ───────────────────
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 1.0, fill=RED)
    add_text(slide, "HISTORIQUE DES SAISIES (20 DERNIÈRES)", 0.3, 0.18, 12, 0.65, size=20, bold=True, color=WHITE)

    if not df_all.empty:
        cols_ppt  = ["nom_resp","categorie","kpi","periode","valeur","unite","date_saisie"]
        col_hdrs  = ["Responsable","Catégorie","KPI","Période","Valeur","Unité","Date"]
        col_ws    = [1.5, 2.0, 3.5, 1.3, 1.1, 0.8, 1.8]
        df_ppt    = df_all[cols_ppt].head(20)

        # Header row
        x = 0.2
        for hdr, w in zip(col_hdrs, col_ws):
            add_rect(slide, x, 1.1, w-0.05, 0.35, fill=DARK)
            add_text(slide, hdr, x+0.05, 1.12, w-0.1, 0.3, size=8, bold=True, color=WHITE)
            x += w

        # Data rows
        for r_i, (_, row) in enumerate(df_ppt.iterrows()):
            bg = LIGHT if r_i % 2 == 0 else WHITE
            x = 0.2
            y_row = 1.5 + r_i * 0.27
            for col_n, w in zip(cols_ppt, col_ws):
                val = str(row[col_n]) if pd.notna(row[col_n]) else ""
                if col_n == "valeur":
                    try: val = f"{float(val):,.2f}"
                    except: pass
                if col_n == "date_saisie":
                    val = val[:10]
                if col_n == "kpi":
                    val = val[:30] + "…" if len(val) > 30 else val
                add_rect(slide, x, y_row, w-0.05, 0.25, fill=bg)
                add_text(slide, val, x+0.04, y_row+0.02, w-0.1, 0.22, size=7, color=DARK)
                x += w

    # ── SLIDE 5 : KPIs manquants ─────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 1.0, fill=RGBColor(0xFF,0xC1,0x07))
    add_text(slide, "⚠  KPIs NON ENCORE SAISIS", 0.3, 0.18, 12, 0.65, size=22, bold=True, color=DARK)

    if not df_all.empty:
        kpis_ok = df_all["kpi"].unique()
        manq = df_ref[~df_ref["kpi"].isin(kpis_ok)]
    else:
        manq = df_ref

    if manq.empty:
        add_text(slide, "✅  Tous les KPIs ont été saisis !", 0.4, 1.4, 12, 1, size=20, bold=True, color=RGBColor(0x2E,0x7D,0x32))
    else:
        y = 1.2
        col = 0
        for _, row in manq.iterrows():
            x = 0.3 + col * 6.5
            add_rect(slide, x, y, 6.3, 0.55, fill=RGBColor(0xFF,0xF8,0xE1))
            add_rect(slide, x, y, 0.1, 0.55, fill=RGBColor(0xFF,0xC1,0x07))
            add_text(slide, row["kpi"], x+0.2, y+0.03, 4.5, 0.28, size=9, bold=True, color=DARK)
            add_text(slide, f"→ {row['nom']}  ·  {row['categorie']}", x+0.2, y+0.28, 5.8, 0.22, size=8, color=GRAY)
            col += 1
            if col == 2:
                col = 0
                y += 0.65
            if y > 6.8:
                break

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── CSS ────────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=IBM+Plex+Sans:wght@300;400;600;700&family=IBM+Plex+Mono&display=swap');
html, body, [class*="css"] { font-family: 'IBM Plex Sans', sans-serif; }
.main { background: #f5f6fa; }
.kpi-card {
    background: white; border-left: 4px solid #cc0000;
    border-radius: 4px; padding: 16px 20px; margin-bottom: 12px;
    box-shadow: 0 1px 4px rgba(0,0,0,0.08);
}
.kpi-val   { font-size: 2rem; font-weight: 700; color: #1a1a2e; font-family: 'IBM Plex Mono'; }
.kpi-label { font-size: 0.78rem; color: #666; text-transform: uppercase; letter-spacing: .08em; }
.section-header {
    font-size: 1.05rem; font-weight: 700; color: #cc0000;
    border-bottom: 2px solid #cc0000; padding-bottom: 4px; margin: 20px 0 12px;
    text-transform: uppercase; letter-spacing: .06em;
}
.alert-box {
    background: #fff8e1; border-left: 4px solid #ffc107;
    padding: 10px 16px; border-radius: 4px; font-size:.88rem; margin-bottom:6px;
}
.badge-admin { background:#1a1a2e; color:white; padding:3px 12px; border-radius:20px; font-size:.8rem; font-weight:600; }
.badge-dg    { background:#7b1fa2; color:white; padding:3px 12px; border-radius:20px; font-size:.8rem; font-weight:600; }
.badge-user  { background:#cc0000; color:white; padding:3px 12px; border-radius:20px; font-size:.8rem; font-weight:600; }
.stat-big { font-size:3rem; font-weight:800; font-family:'IBM Plex Mono'; }
</style>
""", unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════════════════════
# LOGIN
# ═══════════════════════════════════════════════════════════════════════════════
def show_login():
    st.markdown("""
    <div style="text-align:center;margin-top:40px;">
        <img src="https://upload.wikimedia.org/wikipedia/fr/thumb/9/9a/Soci%C3%A9t%C3%A9_G%C3%A9n%C3%A9rale.svg/200px-Soci%C3%A9t%C3%A9_G%C3%A9n%C3%A9rale.svg.png" width="160">
        <h2 style="color:#cc0000;margin-top:16px;">Dashboard KPI – SocGen CI</h2>
        <p style="color:#888;font-size:.9rem;">Saisie &amp; suivi des indicateurs de performance</p>
    </div>
    """, unsafe_allow_html=True)
    _, col, _ = st.columns([1, 1.1, 1])
    with col:
        with st.form("lf"):
            st.markdown("#### 🔐 Connexion")
            li = st.text_input("Identifiant")
            pw = st.text_input("Mot de passe", type="password")
            ok = st.form_submit_button("Se connecter", use_container_width=True, type="primary")
        if ok:
            k = li.strip().lower()
            if k in USERS and USERS[k]["password"] == _hash(pw):
                st.session_state.update({"logged_in": True, "login": k,
                    "role": USERS[k]["role"], "nom_ref": USERS[k]["nom_ref"],
                    "display_name": USERS[k]["display"]})
                log_connexion(k, "LOGIN")
                st.rerun()
            else:
                st.error("Identifiant ou mot de passe incorrect.")
        st.caption("⚠️ Accès restreint – SocGen CI 2026")


if not st.session_state.get("logged_in"):
    show_login()
    st.stop()

CUR_LOGIN   = st.session_state["login"]
CUR_ROLE    = st.session_state["role"]
CUR_NOM_REF = st.session_state["nom_ref"]
CUR_DISPLAY = st.session_state["display_name"]
IS_ADMIN    = CUR_ROLE == "admin"
IS_DG       = CUR_ROLE == "dg"
IS_POWER    = IS_ADMIN or IS_DG   # accès stats + exports

df_ref = load_ref()
df_mes_kpis = df_ref.copy() if IS_POWER else df_ref[df_ref["nom"] == CUR_NOM_REF].copy()


# ═══════════════════════════════════════════════════════════════════════════════
# SIDEBAR
# ═══════════════════════════════════════════════════════════════════════════════
with st.sidebar:
    st.image("https://upload.wikimedia.org/wikipedia/fr/thumb/9/9a/Soci%C3%A9t%C3%A9_G%C3%A9n%C3%A9rale.svg/200px-Soci%C3%A9t%C3%A9_G%C3%A9n%C3%A9rale.svg.png", width=130)
    st.markdown("---")
    badge = f"badge-{CUR_ROLE}"
    label = ROLE_LABELS.get(CUR_ROLE, "Utilisateur")
    st.markdown(f'<span class="{badge}">{label}</span> &nbsp;<b>{CUR_DISPLAY}</b>', unsafe_allow_html=True)
    st.markdown("---")

    if IS_ADMIN:
        pages = ["📝 Dashboard de saisies", "📊 Dashboard Analytique", "📝 Saisie KPI",
                 "🗃️ Historique", "⚙️ Gestion Référentiel", "📋 Référentiel", "🔐 Logs"]
    elif IS_DG:
        pages = ["📊 Dashboard Analytique", "🗃️ Historique", "📋 Référentiel"]
    else:
        pages = ["📝 Dashboard de saisies", "📝 Saisie KPI", "🗃️ Mes saisies", "📋 Référentiel"]

    page = st.radio("Navigation", pages)
    st.markdown("---")
    if st.button("🚪 Déconnecter", use_container_width=True):
        log_connexion(CUR_LOGIN, "LOGOUT")
        for k in list(st.session_state.keys()): del st.session_state[k]
        st.rerun()
    st.caption("SocGen CI · 2026")


# ═══════════════════════════════════════════════════════════════════════════════
# PAGE : DASHBOARD DE SAISIES
# ═══════════════════════════════════════════════════════════════════════════════
if page == "📝 Dashboard de saisies":
    st.markdown(f"# 📝 Dashboard de saisies{' – ' + CUR_DISPLAY if not IS_POWER else ''}")

    df_all      = load_saisies(CUR_LOGIN)
    total_kpis  = len(df_mes_kpis)
    kpis_saisis = df_all["kpi"].nunique() if not df_all.empty else 0
    tx_compl    = round(kpis_saisis / total_kpis * 100) if total_kpis else 0
    nb_saisies  = len(df_all)

    c1,c2,c3,c4 = st.columns(4)
    for col, v, lbl in [(c1,total_kpis,"KPIs assignés"),(c2,kpis_saisis,"KPIs saisis"),
                         (c3,f"{tx_compl}%","Complétion"),(c4,nb_saisies,"Total saisies")]:
        col.markdown(f'<div class="kpi-card"><div class="kpi-label">{lbl}</div><div class="kpi-val">{v}</div></div>', unsafe_allow_html=True)

    st.markdown("---")
    if df_all.empty:
        st.info("Aucune donnée saisie.")
    else:
        cl, cr = st.columns(2)
        with cl:
            st.markdown('<div class="section-header">Complétion par catégorie</div>', unsafe_allow_html=True)
            cat_t = df_mes_kpis.groupby("categorie")["kpi"].count().reset_index(name="total")
            cat_s = df_all.groupby("categorie")["kpi"].nunique().reset_index(name="saisis")
            cat_m = cat_t.merge(cat_s,on="categorie",how="left").fillna(0)
            cat_m["pct"] = (cat_m["saisis"]/cat_m["total"]*100).round(0)
            fig = px.bar(cat_m, x="pct", y="categorie", orientation="h",
                         color="pct", color_continuous_scale=["#ffcdd2","#cc0000"],
                         range_color=[0,100], text="pct", labels={"pct":"%","categorie":""})
            fig.update_traces(texttemplate="%{text:.0f}%", textposition="outside")
            fig.update_layout(showlegend=False, coloraxis_showscale=False,
                              paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                              margin=dict(l=0,r=50,t=10,b=10), height=280,
                              yaxis=dict(autorange="reversed"))
            st.plotly_chart(fig, use_container_width=True)

        with cr:
            if IS_POWER:
                st.markdown('<div class="section-header">Saisies par responsable</div>', unsafe_allow_html=True)
                by_r = df_all.groupby("nom_resp").size().reset_index(name="nb")
                fig2 = px.pie(by_r, values="nb", names="nom_resp",
                              color_discrete_sequence=px.colors.sequential.Reds_r, hole=0.45)
                fig2.update_layout(paper_bgcolor="rgba(0,0,0,0)",
                                   margin=dict(l=0,r=0,t=10,b=10), height=280)
                st.plotly_chart(fig2, use_container_width=True)
            else:
                st.markdown('<div class="section-header">Mes dernières saisies</div>', unsafe_allow_html=True)
                st.dataframe(df_all[["kpi","periode","valeur","unite","date_saisie"]].head(8)
                    .rename(columns={"kpi":"KPI","periode":"Période","valeur":"Valeur","unite":"Unité","date_saisie":"Date"}),
                    use_container_width=True, hide_index=True)

        # Évolution
        st.markdown('<div class="section-header">Évolution des saisies</div>', unsafe_allow_html=True)
        df_all["dt"] = pd.to_datetime(df_all["date_saisie"])
        df_tl = df_all.groupby(df_all["dt"].dt.date).size().reset_index(name="nb")
        df_tl.columns = ["date","nb"]
        fig3 = px.area(df_tl, x="date", y="nb", color_discrete_sequence=["#cc0000"], labels={"nb":"Saisies","date":""})
        fig3.update_layout(paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                           margin=dict(l=0,r=0,t=10,b=10), height=190)
        st.plotly_chart(fig3, use_container_width=True)

        # KPIs manquants avec nom du responsable
        kpis_ok = df_all["kpi"].unique()
        manquants = df_mes_kpis[~df_mes_kpis["kpi"].isin(kpis_ok)]
        if not manquants.empty:
            st.markdown('<div class="section-header">⚠️ KPIs non encore saisis</div>', unsafe_allow_html=True)
            for _, row in manquants.iterrows():
                st.markdown(
                    f'<div class="alert-box">👤 <b>{row["nom"]}</b> &nbsp;·&nbsp; '
                    f'<b>{row["categorie"]}</b> — {row["kpi"]} '
                    f'<span style="color:#888">· {row["periodicite"]}</span></div>',
                    unsafe_allow_html=True
                )


# ═══════════════════════════════════════════════════════════════════════════════
# PAGE : DASHBOARD ANALYTIQUE (DG + Admin)
# ═══════════════════════════════════════════════════════════════════════════════
elif page == "📊 Dashboard Analytique":
    if not IS_POWER:
        st.error("Accès réservé à la Direction Générale et à l'Admin.")
        st.stop()

    st.markdown("# 📊 Dashboard Analytique – Direction Générale")

    df_all = load_saisies()  # toutes les données
    df_ref_full = load_ref()

    if df_all.empty:
        st.info("Aucune donnée disponible dans l'historique.")
        st.stop()

    # ── Filtres ──
    with st.expander("🔍 Filtres", expanded=True):
        fc1, fc2, fc3 = st.columns(3)
        with fc1:
            cats = df_all["categorie"].unique().tolist()
            f_cats = st.multiselect("Catégorie", cats, default=cats)
        with fc2:
            resps = df_all["nom_resp"].unique().tolist()
            f_resps = st.multiselect("Responsable", resps, default=resps)
        with fc3:
            periodes = sorted(df_all["periode"].unique().tolist(), reverse=True)
            f_per = st.multiselect("Période", periodes, default=periodes[:3] if len(periodes)>=3 else periodes)

    df_f = df_all[
        df_all["categorie"].isin(f_cats) &
        df_all["nom_resp"].isin(f_resps) &
        (df_all["periode"].isin(f_per) if f_per else True)
    ]

    # ── KPIs stats ──
    c1,c2,c3,c4,c5 = st.columns(5)
    for col, v, lbl in [
        (c1, len(df_f), "Saisies filtrées"),
        (c2, df_f["nom_resp"].nunique(), "Responsables actifs"),
        (c3, df_f["categorie"].nunique(), "Catégories couvertes"),
        (c4, df_f["kpi"].nunique(), "KPIs distincts"),
        (c5, f"{df_f['valeur'].mean():.2f}" if not df_f.empty else "—", "Valeur moyenne"),
    ]:
        col.markdown(f'<div class="kpi-card"><div class="kpi-label">{lbl}</div><div class="kpi-val" style="font-size:1.6rem">{v}</div></div>', unsafe_allow_html=True)

    st.markdown("---")

    row1_l, row1_r = st.columns(2)

    # Volume par catégorie
    with row1_l:
        st.markdown('<div class="section-header">Saisies par catégorie</div>', unsafe_allow_html=True)
        by_cat = df_f.groupby("categorie").size().reset_index(name="nb").sort_values("nb", ascending=True)
        fig = px.bar(by_cat, x="nb", y="categorie", orientation="h",
                     color="nb", color_continuous_scale=["#ffcdd2","#cc0000"],
                     text="nb", labels={"nb":"Saisies","categorie":""})
        fig.update_traces(textposition="outside")
        fig.update_layout(showlegend=False, coloraxis_showscale=False,
                          paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                          margin=dict(l=0,r=40,t=10,b=10), height=300)
        st.plotly_chart(fig, use_container_width=True)

    # Saisies par responsable
    with row1_r:
        st.markdown('<div class="section-header">Saisies par responsable</div>', unsafe_allow_html=True)
        by_r = df_f.groupby("nom_resp").size().reset_index(name="nb").sort_values("nb", ascending=False)
        fig2 = px.bar(by_r, x="nom_resp", y="nb",
                      color="nb", color_continuous_scale=["#ffcdd2","#cc0000"],
                      text="nb", labels={"nb":"Saisies","nom_resp":"Responsable"})
        fig2.update_traces(textposition="outside")
        fig2.update_layout(showlegend=False, coloraxis_showscale=False,
                           paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                           margin=dict(l=0,r=0,t=10,b=10), height=300,
                           xaxis_tickangle=-30)
        st.plotly_chart(fig2, use_container_width=True)

    row2_l, row2_r = st.columns(2)

    # Évolution par période
    with row2_l:
        st.markdown('<div class="section-header">Volume de saisies par période</div>', unsafe_allow_html=True)
        by_per = df_f.groupby("periode").size().reset_index(name="nb")
        fig3 = px.bar(by_per, x="periode", y="nb",
                      color_discrete_sequence=["#cc0000"],
                      text="nb", labels={"nb":"Saisies","periode":"Période"})
        fig3.update_traces(textposition="outside")
        fig3.update_layout(paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                           margin=dict(l=0,r=0,t=10,b=10), height=280)
        st.plotly_chart(fig3, use_container_width=True)

    # Distribution des valeurs
    with row2_r:
        st.markdown('<div class="section-header">Distribution des valeurs saisies</div>', unsafe_allow_html=True)
        fig4 = px.box(df_f, x="categorie", y="valeur",
                      color_discrete_sequence=["#cc0000"],
                      labels={"valeur":"Valeur","categorie":"Catégorie"})
        fig4.update_layout(paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                           margin=dict(l=0,r=0,t=10,b=10), height=280,
                           xaxis_tickangle=-20)
        st.plotly_chart(fig4, use_container_width=True)

    # Tableau pivot valeurs moyennes par KPI × période
    st.markdown('<div class="section-header">Tableau pivot : Valeur moyenne par KPI et période</div>', unsafe_allow_html=True)
    if not df_f.empty and len(f_per) > 0:
        pivot = df_f.pivot_table(values="valeur", index="kpi", columns="periode", aggfunc="mean").round(2)
        st.dataframe(pivot, use_container_width=True, height=300)

    # ── Exports ──
    st.markdown("---")
    st.markdown('<div class="section-header">📥 Exports</div>', unsafe_allow_html=True)
    exp1, exp2 = st.columns(2)

    with exp1:
        xlsx_bytes = to_excel_bytes(
            df_f.rename(columns={"id":"ID","login":"Login","nom_resp":"Responsable",
                                  "categorie":"Catégorie","kpi":"KPI","periodicite":"Périodicité",
                                  "periode":"Période","valeur":"Valeur","unite":"Unité",
                                  "commentaire":"Commentaire","date_saisie":"Date saisie"})
        )
        st.download_button(
            "⬇️ Télécharger l'historique (.xlsx)",
            data=xlsx_bytes,
            file_name=f"kpi_historique_{datetime.now().strftime('%Y%m%d')}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True
        )

    with exp2:
        pptx_bytes = build_pptx(df_all, df_ref_full)
        st.download_button(
            "⬇️ Télécharger le dashboard (.pptx)",
            data=pptx_bytes,
            file_name=f"kpi_dashboard_{datetime.now().strftime('%Y%m%d')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )


# ═══════════════════════════════════════════════════════════════════════════════
# PAGE : SAISIE KPI
# ═══════════════════════════════════════════════════════════════════════════════
elif page == "📝 Saisie KPI":
    st.markdown("# 📝 Formulaire de saisie KPI")
    if df_mes_kpis.empty:
        st.warning("Aucun KPI assigné à votre compte.")
        st.stop()

    cats = df_mes_kpis["categorie"].unique().tolist()
    with st.form("fs", clear_on_submit=True):
        c1, c2 = st.columns(2)
        with c1:
            st.markdown("**Identification**")
            categorie  = st.selectbox("Catégorie *", cats)
            kpis_liste = df_mes_kpis[df_mes_kpis["categorie"] == categorie]["kpi"].tolist()
            kpi        = st.selectbox("KPI *", kpis_liste)
        with c2:
            st.markdown("**Données**")
            periode     = st.text_input("Période *", value=datetime.now().strftime("%B %Y"))
            valeur      = st.number_input("Valeur *", value=0.0, format="%.4f")
            unite       = st.text_input("Unité", placeholder="%, FCFA, nombre…")
            commentaire = st.text_area("Commentaire", height=80)

        info_kpi = df_mes_kpis[df_mes_kpis["kpi"] == kpi].iloc[0] if kpi else None
        if info_kpi is not None:
            st.info(f"📧 **{info_kpi['nom']}** ({info_kpi['email']}) · Périodicité : **{info_kpi['periodicite']}**")

        if st.form_submit_button("✅ Enregistrer", use_container_width=True, type="primary"):
            if not periode.strip():
                st.error("La période est obligatoire.")
            else:
                insert_saisie({
                    "login": CUR_LOGIN, "nom_resp": CUR_NOM_REF or CUR_DISPLAY,
                    "categorie": categorie, "kpi": kpi,
                    "periodicite": info_kpi["periodicite"] if info_kpi is not None else "",
                    "periode": periode.strip(), "valeur": valeur,
                    "unite": unite, "commentaire": commentaire,
                })
                st.success(f"✅ **{kpi}** – période **{periode}** enregistré !")


# ═══════════════════════════════════════════════════════════════════════════════
# PAGE : HISTORIQUE
# ═══════════════════════════════════════════════════════════════════════════════
elif page in ("🗃️ Historique", "🗃️ Mes saisies"):
    title = "🗃️ Historique global" if IS_POWER else "🗃️ Mes saisies"
    st.markdown(f"# {title}")
    df_h = load_saisies(CUR_LOGIN)

    if df_h.empty:
        st.info("Aucune donnée disponible.")
    else:
        c1, c2, c3 = st.columns(3)
        cats_u = df_h["categorie"].unique().tolist()
        with c1: f_cat = st.multiselect("Catégorie", cats_u, default=cats_u)
        with c2: f_per = st.text_input("Période (texte libre)")
        f_resp = None
        with c3:
            if IS_POWER:
                r_u = df_h["nom_resp"].unique().tolist()
                f_resp = st.multiselect("Responsable", r_u, default=r_u)

        df_f = df_h[df_h["categorie"].isin(f_cat)]
        if f_per: df_f = df_f[df_f["periode"].str.contains(f_per, case=False, na=False)]
        if IS_POWER and f_resp: df_f = df_f[df_f["nom_resp"].isin(f_resp)]

        cols_show = (["id","nom_resp","categorie","kpi","periodicite","periode","valeur","unite","commentaire","date_saisie"]
                     if IS_POWER else ["id","categorie","kpi","periodicite","periode","valeur","unite","commentaire","date_saisie"])
        rename_map = {"id":"ID","nom_resp":"Responsable","categorie":"Catégorie","kpi":"KPI",
                      "periodicite":"Périodicité","periode":"Période","valeur":"Valeur",
                      "unite":"Unité","commentaire":"Commentaire","date_saisie":"Date saisie"}
        df_display = df_f[cols_show].rename(columns=rename_map)
        st.dataframe(df_display, use_container_width=True, height=400)
        st.caption(f"{len(df_f)} enregistrement(s)")

        # Exports
        exp_c, exp_x = st.columns(2)
        with exp_c:
            st.download_button("⬇️ Exporter CSV", data=df_f.to_csv(index=False).encode("utf-8"),
                               file_name=f"kpi_{CUR_LOGIN}.csv", mime="text/csv")
        if IS_POWER:
            with exp_x:
                xlsx_bytes = to_excel_bytes(df_display)
                st.download_button("⬇️ Exporter Excel (.xlsx)", data=xlsx_bytes,
                                   file_name=f"kpi_historique_{datetime.now().strftime('%Y%m%d')}.xlsx",
                                   mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

        with st.expander("🗑️ Supprimer un enregistrement"):
            del_id = st.number_input("ID à supprimer", min_value=1, step=1)
            if st.button("Supprimer", type="secondary"):
                delete_saisie(int(del_id), CUR_LOGIN)
                st.success(f"Ligne #{del_id} supprimée.")
                st.rerun()


# ═══════════════════════════════════════════════════════════════════════════════
# PAGE : GESTION RÉFÉRENTIEL (Admin)
# ═══════════════════════════════════════════════════════════════════════════════
elif page == "⚙️ Gestion Référentiel":
    if not IS_ADMIN:
        st.error("Accès réservé à l'administrateur.")
        st.stop()

    st.markdown("# ⚙️ Gestion du Référentiel KPI")
    st.info("Ajoutez de nouvelles catégories ou de nouveaux KPIs. Les modifications sont prises en compte immédiatement.")

    # ── Ajouter un KPI ──
    st.markdown('<div class="section-header">Ajouter un KPI</div>', unsafe_allow_html=True)
    df_ref_all = load_ref()
    cats_existantes = sorted(df_ref_all["categorie"].unique().tolist())
    resps_existants = sorted(df_ref_all["nom"].dropna().unique().tolist())

    with st.form("add_kpi_form", clear_on_submit=True):
        ac1, ac2 = st.columns(2)
        with ac1:
            choix_cat = st.radio("Catégorie", ["Choisir existante", "Créer nouvelle"], horizontal=True)
            if choix_cat == "Choisir existante":
                new_cat = st.selectbox("Catégorie *", cats_existantes)
            else:
                new_cat = st.text_input("Nom de la nouvelle catégorie *")

            new_kpi = st.text_input("Libellé du KPI *")
            new_per = st.selectbox("Périodicité *", ["Mensuelle", "Trimestrielle", "Semestrielle", "Annuelle"])

        with ac2:
            choix_resp = st.radio("Responsable", ["Choisir existant", "Nouveau responsable"], horizontal=True)
            if choix_resp == "Choisir existant":
                new_nom = st.selectbox("Responsable *", resps_existants)
                # Pré-remplir l'email si dispo
                email_auto = df_ref_all[df_ref_all["nom"] == new_nom]["email"].values
                new_email = st.text_input("Email *", value=email_auto[0] if len(email_auto) > 0 else "")
            else:
                new_nom   = st.text_input("Nom du responsable *")
                new_email = st.text_input("Email du responsable *")

        submitted = st.form_submit_button("➕ Ajouter le KPI", use_container_width=True, type="primary")
        if submitted:
            if not new_cat.strip() or not new_kpi.strip() or not new_nom.strip():
                st.error("Catégorie, KPI et Responsable sont obligatoires.")
            else:
                add_kpi_ref(new_cat.strip(), new_kpi.strip(), new_per, new_email.strip(), new_nom.strip())
                st.success(f"✅ KPI **{new_kpi}** ajouté dans la catégorie **{new_cat}**.")
                st.rerun()

    # ── Liste et suppression ──
    st.markdown('<div class="section-header">KPIs actifs dans le référentiel</div>', unsafe_allow_html=True)
    df_ref_cur = load_ref()
    for cat in df_ref_cur["categorie"].unique():
        with st.expander(f"📁 {cat} ({len(df_ref_cur[df_ref_cur['categorie']==cat])} KPIs)"):
            sub = df_ref_cur[df_ref_cur["categorie"] == cat][["id","kpi","periodicite","nom","email"]]
            st.dataframe(sub.rename(columns={"id":"ID","kpi":"KPI","periodicite":"Périodicité",
                                              "nom":"Responsable","email":"Email"}),
                         use_container_width=True, hide_index=True)
            del_kpi_id = st.number_input(f"ID à désactiver (catégorie {cat})", min_value=1, step=1, key=f"del_{cat}")
            if st.button(f"🗑️ Désactiver ce KPI", key=f"btn_{cat}", type="secondary"):
                delete_kpi_ref(int(del_kpi_id))
                st.success(f"KPI #{del_kpi_id} désactivé.")
                st.rerun()


# ═══════════════════════════════════════════════════════════════════════════════
# PAGE : RÉFÉRENTIEL
# ═══════════════════════════════════════════════════════════════════════════════
elif page == "📋 Référentiel":
    st.markdown("# 📋 Référentiel des KPIs")
    if not IS_POWER:
        st.info(f"Affichage filtré : KPIs assignés à **{CUR_DISPLAY}**")

    search = st.text_input("🔍 Rechercher…")
    df_show = df_mes_kpis.copy()
    if search:
        mask = df_show.apply(lambda r: r.astype(str).str.contains(search, case=False).any(), axis=1)
        df_show = df_show[mask]

    for cat in df_show["categorie"].unique():
        st.markdown(f'<div class="section-header">{cat}</div>', unsafe_allow_html=True)
        sub = (df_show[df_show["categorie"] == cat][["kpi","periodicite","nom","email"]]
               .rename(columns={"kpi":"KPI","periodicite":"Périodicité","nom":"Responsable","email":"Email"}))
        st.dataframe(sub, use_container_width=True, hide_index=True)


# ═══════════════════════════════════════════════════════════════════════════════
# PAGE : LOGS (Admin)
# ═══════════════════════════════════════════════════════════════════════════════
elif page == "🔐 Logs":
    if not IS_ADMIN:
        st.error("Accès réservé à l'administrateur.")
        st.stop()
    st.markdown("# 🔐 Logs de connexion")
    df_logs = load_logs()
    if df_logs.empty:
        st.info("Aucun log.")
    else:
        st.dataframe(df_logs.rename(columns={"id":"ID","login":"Login","action":"Action","timestamp":"Horodatage"}),
                     use_container_width=True, height=400)
        st.download_button("⬇️ Exporter logs CSV", data=df_logs.to_csv(index=False).encode(),
                           file_name="logs.csv", mime="text/csv")
